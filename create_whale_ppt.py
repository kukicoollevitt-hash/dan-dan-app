#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 프레젠테이션 생성 (16:9 비율)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 색상 정의
WHALE_BLUE = RGBColor(30, 60, 114)      # 진한 남색
LIGHT_BLUE = RGBColor(66, 133, 244)     # 밝은 파랑
ACCENT_ORANGE = RGBColor(255, 152, 0)   # 주황
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(50, 50, 50)
LIGHT_GRAY = RGBColor(245, 245, 245)

def add_title_slide(prs, title, subtitle=""):
    """표지 슬라이드"""
    slide_layout = prs.slide_layouts[6]  # 빈 슬라이드
    slide = prs.slides.add_slide(slide_layout)

    # 배경색
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = WHALE_BLUE
    background.line.fill.background()

    # 고래 이모지
    whale = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(2), Inches(1.5))
    whale_tf = whale.text_frame
    whale_p = whale_tf.paragraphs[0]
    whale_p.text = "🐋"
    whale_p.font.size = Pt(100)
    whale_p.alignment = PP_ALIGN.CENTER

    # 제목
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(11.333), Inches(1.2))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    title_p.alignment = PP_ALIGN.CENTER

    # 부제목
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.333), Inches(0.8))
        sub_tf = sub_box.text_frame
        sub_p = sub_tf.paragraphs[0]
        sub_p.text = subtitle
        sub_p.font.size = Pt(28)
        sub_p.font.color.rgb = RGBColor(200, 220, 255)
        sub_p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, content_lines, emoji=""):
    """일반 콘텐츠 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 흰색 배경
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = WHITE
    background.line.fill.background()

    # 상단 파란 바
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = WHALE_BLUE
    top_bar.line.fill.background()

    # 제목
    title_text = f"{emoji} {title}" if emoji else title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = title_text
    title_p.font.size = Pt(36)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE

    # 콘텐츠
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(12), Inches(5.5))
    content_tf = content_box.text_frame
    content_tf.word_wrap = True

    for i, line in enumerate(content_lines):
        if i == 0:
            p = content_tf.paragraphs[0]
        else:
            p = content_tf.add_paragraph()
        p.text = line
        p.font.size = Pt(22)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(12)

    return slide

def add_table_slide(prs, title, headers, rows, emoji=""):
    """테이블 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 흰색 배경
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = WHITE
    background.line.fill.background()

    # 상단 파란 바
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = WHALE_BLUE
    top_bar.line.fill.background()

    # 제목
    title_text = f"{emoji} {title}" if emoji else title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = title_text
    title_p.font.size = Pt(36)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE

    # 테이블
    num_rows = len(rows) + 1
    num_cols = len(headers)

    table_width = Inches(11)
    table_height = Inches(0.6 * num_rows)
    left = Inches(1.166)
    top = Inches(1.8)

    table = slide.shapes.add_table(num_rows, num_cols, left, top, table_width, table_height).table

    # 헤더
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHALE_BLUE
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.color.rgb = WHITE
        para.font.size = Pt(18)
        para.alignment = PP_ALIGN.CENTER

    # 데이터
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY if row_idx % 2 == 0 else WHITE
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(16)
            para.font.color.rgb = DARK_GRAY
            para.alignment = PP_ALIGN.CENTER

    return slide

def add_three_column_slide(prs, title, col1, col2, col3, emoji=""):
    """3열 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 흰색 배경
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = WHITE
    background.line.fill.background()

    # 상단 파란 바
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = WHALE_BLUE
    top_bar.line.fill.background()

    # 제목
    title_text = f"{emoji} {title}" if emoji else title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = title_text
    title_p.font.size = Pt(36)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE

    # 3개 컬럼
    columns = [col1, col2, col3]
    colors = [LIGHT_BLUE, ACCENT_ORANGE, RGBColor(76, 175, 80)]

    for i, (col_title, col_content) in enumerate(columns):
        left = Inches(0.8 + i * 4.1)

        # 박스
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.6), Inches(3.8), Inches(5.3))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = colors[i]
        box.line.width = Pt(3)

        # 컬럼 제목
        col_title_box = slide.shapes.add_textbox(left + Inches(0.2), Inches(1.8), Inches(3.4), Inches(0.6))
        col_tf = col_title_box.text_frame
        col_p = col_tf.paragraphs[0]
        col_p.text = col_title
        col_p.font.size = Pt(22)
        col_p.font.bold = True
        col_p.font.color.rgb = colors[i]
        col_p.alignment = PP_ALIGN.CENTER

        # 컬럼 내용
        content_box = slide.shapes.add_textbox(left + Inches(0.2), Inches(2.5), Inches(3.4), Inches(4))
        content_tf = content_box.text_frame
        content_tf.word_wrap = True

        for j, line in enumerate(col_content):
            if j == 0:
                p = content_tf.paragraphs[0]
            else:
                p = content_tf.add_paragraph()
            p.text = f"• {line}"
            p.font.size = Pt(16)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(10)

    return slide

def add_number_slide(prs, title, numbers, emoji=""):
    """숫자 강조 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 배경
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = WHALE_BLUE
    background.line.fill.background()

    # 제목
    title_text = f"{emoji} {title}" if emoji else title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = title_text
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    title_p.alignment = PP_ALIGN.CENTER

    # 숫자들
    num_items = len(numbers)
    for i, (num, label) in enumerate(numbers):
        left = Inches(1 + i * (11 / num_items))

        # 숫자
        num_box = slide.shapes.add_textbox(left, Inches(2.5), Inches(3), Inches(1.5))
        num_tf = num_box.text_frame
        num_p = num_tf.paragraphs[0]
        num_p.text = num
        num_p.font.size = Pt(72)
        num_p.font.bold = True
        num_p.font.color.rgb = ACCENT_ORANGE
        num_p.alignment = PP_ALIGN.CENTER

        # 라벨
        label_box = slide.shapes.add_textbox(left, Inches(4.2), Inches(3), Inches(0.8))
        label_tf = label_box.text_frame
        label_p = label_tf.paragraphs[0]
        label_p.text = label
        label_p.font.size = Pt(24)
        label_p.font.color.rgb = WHITE
        label_p.alignment = PP_ALIGN.CENTER

    return slide

def add_flow_slide(prs, title, steps, emoji=""):
    """흐름도 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 흰색 배경
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = WHITE
    background.line.fill.background()

    # 상단 파란 바
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = WHALE_BLUE
    top_bar.line.fill.background()

    # 제목
    title_text = f"{emoji} {title}" if emoji else title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = title_text
    title_p.font.size = Pt(36)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE

    # 단계들
    num_steps = len(steps)
    box_width = Inches(2.2)
    spacing = (Inches(12) - box_width * num_steps) / (num_steps + 1)

    for i, (step_num, step_text) in enumerate(steps):
        left = spacing + i * (box_width + spacing)

        # 원형 번호
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.7), Inches(2), Inches(0.8), Inches(0.8))
        circle.fill.solid()
        circle.fill.fore_color.rgb = WHALE_BLUE
        circle.line.fill.background()

        circle_text = slide.shapes.add_textbox(left + Inches(0.7), Inches(2.1), Inches(0.8), Inches(0.7))
        circle_tf = circle_text.text_frame
        circle_p = circle_tf.paragraphs[0]
        circle_p.text = step_num
        circle_p.font.size = Pt(28)
        circle_p.font.bold = True
        circle_p.font.color.rgb = WHITE
        circle_p.alignment = PP_ALIGN.CENTER

        # 박스
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(3), box_width, Inches(2))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = WHALE_BLUE
        box.line.width = Pt(2)

        # 텍스트
        text_box = slide.shapes.add_textbox(left + Inches(0.1), Inches(3.3), Inches(2), Inches(1.5))
        text_tf = text_box.text_frame
        text_tf.word_wrap = True
        text_p = text_tf.paragraphs[0]
        text_p.text = step_text
        text_p.font.size = Pt(16)
        text_p.font.color.rgb = DARK_GRAY
        text_p.alignment = PP_ALIGN.CENTER

        # 화살표 (마지막 아님)
        if i < num_steps - 1:
            arrow_left = left + box_width + Inches(0.1)
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_left, Inches(3.8), spacing - Inches(0.2), Inches(0.4))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT_ORANGE
            arrow.line.fill.background()

    return slide

def add_cta_slide(prs, title, subtitle, contact):
    """CTA 슬라이드"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 배경
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = WHALE_BLUE
    background.line.fill.background()

    # 고래 이모지
    whale = slide.shapes.add_textbox(Inches(5.5), Inches(1), Inches(2), Inches(1.5))
    whale_tf = whale.text_frame
    whale_p = whale_tf.paragraphs[0]
    whale_p.text = "🐋"
    whale_p.font.size = Pt(80)
    whale_p.alignment = PP_ALIGN.CENTER

    # 제목
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11.333), Inches(1))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    title_p.alignment = PP_ALIGN.CENTER

    # 부제목
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.333), Inches(0.8))
    sub_tf = sub_box.text_frame
    sub_p = sub_tf.paragraphs[0]
    sub_p.text = subtitle
    sub_p.font.size = Pt(24)
    sub_p.font.color.rgb = RGBColor(200, 220, 255)
    sub_p.alignment = PP_ALIGN.CENTER

    # 연락처 박스
    contact_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(5.2), Inches(5.333), Inches(1.2))
    contact_box.fill.solid()
    contact_box.fill.fore_color.rgb = ACCENT_ORANGE
    contact_box.line.fill.background()

    contact_text = slide.shapes.add_textbox(Inches(4), Inches(5.4), Inches(5.333), Inches(0.9))
    contact_tf = contact_text.text_frame
    contact_p = contact_tf.paragraphs[0]
    contact_p.text = contact
    contact_p.font.size = Pt(28)
    contact_p.font.bold = True
    contact_p.font.color.rgb = WHITE
    contact_p.alignment = PP_ALIGN.CENTER

    return slide

# ========== 슬라이드 생성 ==========

# 1. 표지
add_title_slide(prs, "고래도서관", '"읽는 즐거움, 자라는 실력"\n\n초3 ~ 중2 대상 인터랙티브 독서 학습 플랫폼')

# 2. 문제 제기
add_content_slide(prs, "요즘 아이들, 책 읽기 어려워합니다", [
    "",
    "😟  유튜브, 게임에 익숙한 세대",
    "",
    "😴  긴 글 읽기 집중력 부족",
    "",
    "📖  독서 = 지루하다는 인식",
    "",
    "🏫  학원에서 독서 프로그램 차별화 어려움",
    "",
    "",
    "➡️  해결책이 필요합니다"
], "📌")

# 3. 솔루션
add_content_slide(prs, "고래도서관이 답입니다", [
    "",
    "✅  게임처럼 재미있는 시즌제 스토리",
    "",
    "✅  관문 퀴즈로 성취감 UP",
    "",
    "✅  교과 연계로 학습 효과까지",
    "",
    "✅  감상문 작성으로 글쓰기 능력 향상",
    "",
    "",
    '"놀면서 읽고, 읽으면서 배운다"'
], "💡")

# 4. 판타지 시리즈
add_table_slide(prs, "판타지 시리즈 4종",
    ["시리즈", "테마", "구성"],
    [
        ["🐋 푸른 고래섬의 비밀 모험단", "바다 모험", "3시즌 × 10화"],
        ["🌿 깊은 정글의 잊힌 도시", "정글 탐험", "3시즌 × 10화"],
        ["📚 살아 있는 신비의 도서관", "마법 도서관", "3시즌 × 10화"],
        ["⏰ 시간도서관", "시간 여행", "3시즌 × 10화"],
    ],
    "📚"
)

# 5. 교과 창작도서
add_table_slide(prs, "교과기반 창작도서 312종",
    ["시즌", "대상 학년", "도서 수", "총 화수"],
    [
        ["시즌1", "초3-4학년", "104종", "1,040화"],
        ["시즌2", "초5-6학년", "104종", "1,040화"],
        ["시즌3", "중1-2", "104종", "1,040화"],
    ],
    "📖"
)

# 6. 3단계 학습 시스템
add_flow_slide(prs, "3단계 학습 시스템", [
    ("1", "📖\n스토리 읽기\n\n10화 완독\n몰입 독서"),
    ("2", "🎮\n관문 퀴즈\n\n이해도 확인\n성취감 부여"),
    ("3", "✍️\n감상문 작성\n\n글쓰기 능력\n사고력 향상"),
], "🎯")

# 7. 시즌제 구조
add_content_slide(prs, "시즌제로 성취감 극대화", [
    "",
    "시즌1 (10화)  →  🔓 관문퀴즈  →  시즌2 (10화)  →  🔓 관문퀴즈  →  시즌3 (10화)  →  🎉 완료!",
    "",
    "",
    "🎮  다음 시즌 '해금' 시스템",
    "",
    "🔥  게임처럼 도전 욕구 자극",
    "",
    "📚  끝까지 읽게 만드는 구조"
], "🏆")

# 8. 주요 기능
add_table_slide(prs, "핵심 기능",
    ["기능", "설명"],
    [
        ["🔍 검색", "도서/시리즈 빠른 검색"],
        ["💾 자동저장", "읽던 위치 자동 기억"],
        ["📊 진행바", "실시간 독서 진행도 표시"],
        ["📄 PDF", "감상문 PDF 다운로드"],
        ["📱 카톡공유", "감상문 간편 공유"],
        ["🎨 UI/UX", "몰입형 애니메이션"],
    ],
    "⚙️"
)

# 9. 숫자로 보는 고래도서관
add_number_slide(prs, "숫자로 보는 고래도서관", [
    ("4종", "판타지 시리즈"),
    ("312종", "교과 창작도서"),
    ("3,240화", "총 콘텐츠"),
    ("6학년", "초3~중2 커버"),
], "📊")

# 10. 타겟별 가치
add_three_column_slide(prs, "누구에게 좋은가요?",
    ("👧 학생", [
        "게임처럼 재미있어요",
        "관문 통과할 때 뿌듯해요",
        "어디서든 읽을 수 있어요",
        "친구들과 공유해요"
    ]),
    ("👩 학부모", [
        "학습 진행도 확인 가능",
        "감상문 PDF 포트폴리오",
        "교과 연계로 학교 공부 도움",
        "안전한 콘텐츠"
    ]),
    ("🏫 학원", [
        "독서 프로그램 차별화",
        "방대한 콘텐츠로 장기 운영",
        "학부모 만족도 향상",
        "신규 회원 유치 효과"
    ]),
    "👥"
)

# 11. 경쟁 우위
add_table_slide(prs, "왜 고래도서관인가?",
    ["항목", "기존 독서 프로그램", "고래도서관"],
    [
        ["콘텐츠 양", "적음", "3,240화 ⭐"],
        ["재미 요소", "단순 텍스트", "시즌제 + 관문 ⭐"],
        ["교과 연계", "없거나 약함", "312종 연계 ⭐"],
        ["감상문", "수기 작성", "PDF 자동생성 ⭐"],
        ["진행 관리", "어려움", "자동 추적 ⭐"],
        ["학생 몰입", "낮음", "게임형 구조 ⭐"],
    ],
    "🥇"
)

# 12. 사용자 여정
add_flow_slide(prs, "학습 흐름", [
    ("1", "🔐\n로그인"),
    ("2", "📚\n시리즈 선택"),
    ("3", "📖\n스토리 읽기"),
    ("4", "🎮\n관문 퀴즈"),
    ("5", "✍️\n감상문 작성"),
], "🚀")

# 13. 도입 효과
add_content_slide(prs, "기대 효과", [
    "",
    "📚  독서량          ⬆️  3배 이상 증가",
    "",
    "✍️  글쓰기 능력     ⬆️  감상문 습관화",
    "",
    "🎯  집중력          ⬆️  시즌 완독률 UP",
    "",
    "😊  만족도          ⬆️  재등록률 향상",
    "",
    "",
    '"재미있으니까 스스로 읽습니다"'
], "📈")

# 14. CTA
add_cta_slide(prs,
    "고래도서관과 함께하세요",
    "3,240화 콘텐츠  |  초3 ~ 중2 전 학년  |  교과 연계 + 판타지",
    "📞 도입 문의"
)

# 저장
output_path = "/Users/dandan/Desktop/brainmoon_academy0329/고래도서관_판매PPT.pptx"
prs.save(output_path)
print(f"✅ PPT 생성 완료: {output_path}")
